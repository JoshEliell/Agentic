from django.shortcuts import render
from django.http import HttpResponse
from .forms import FileFieldForm
import json
# utils.py
import tempfile
from pptx import Presentation
import fitz  # PyMuPDF
import subprocess
from langchain_huggingface.embeddings import HuggingFaceEmbeddings
from langchain_community.vectorstores import FAISS
from langchain_ollama import OllamaLLM
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain.schema import Document
from langchain.chains import RetrievalQA
#Agentics
from langchain.prompts import PromptTemplate
#Docs
from docx import Document as DocxDocument
from docx.shared import Inches
from docx.enum.text import WD_ALIGN_PARAGRAPH
from datetime import datetime
from io import BytesIO # ¡Importa BytesIO!
import re   
#Youtube videos
from youtube_transcript_api import YouTubeTranscriptApi
from urllib.parse import urlparse, parse_qs

# Create your views here.
def home(request):
    resultados = None
    preguntas = None
    
    if request.method == 'POST':
        form = FileFieldForm(request.POST, request.FILES)
        if form.is_valid():
            archivos = request.FILES.getlist('Files')
            youtube_link = request.POST.get('youtube_link', '').strip()
            texto_total = ""
            
            for archivo in archivos:
                texto_total += extraer_texto(archivo) + "\n"

            # Extract text from YouTube video
            if youtube_link:
                youtube_text = extract_youtube_text(youtube_link)
                total_text += "\n" + youtube_text

            if not texto_total.strip():
                print("No content to process")
                return render(request, 'dashboard/dashboard.html', {'form': form})
            
            resultados_raw = extract_insights_with_rag(texto_total)
            json_string_data = resultados_raw.get("result", "{}")

            try:
                resultados = json.loads(json_string_data)
            except json.JSONDecodeError as e:
                print(f"Error al decodificar JSON del LLM: {e}")
                resultados = {}

            # Asignación segura
            campaign_objectives = resultados.get("campaign_objectives", [])
            product_specifications = resultados.get("product_specifications", [])
            audience_data = resultados.get("audience_segments", [])
            content_type = resultados.get("content_type", [])
            tone_style = resultados.get("tone_and_style", [])
            historical_context = resultados.get("historical_context", [])
            business_insights = resultados.get("business_insights", [])
            stakeholder_comments = resultados.get("stakeholder_comments", [])
            content_tags = resultados.get("content_tags", [])
            delivery_tools = resultados.get("delivery_platforms", [])

            # Lógica de agentes
            agent_1 = agentic_objetive(campaign_objectives)
            print('20%')
            agent_2 = agentic_product_analysis(product_specifications, audience_data, content_type, tone_style)
            agent_3 = agentic_background(historical_context)
            print('60%')
            agent_4 = agentic_audience_profiling(audience_data, agent_2)
            agent_5 = agentic_insights(product_specifications, campaign_objectives, agent_4, agent_2)
            print('80%')
            final_brief = agentic_brief_generator(product_specifications, agent_4, agent_1, agent_5, agent_3)
            print('90%')
            preguntas = final_brief

            # Generar el archivo Word
            docx_buffer = generate_brief_docx_in_memory(preguntas)
            response = HttpResponse(docx_buffer.getvalue(), content_type='application/vnd.openxmlformats-officedocument.wordprocessingml.document')
            response['Content-Disposition'] = f'attachment; filename="brief_generate_{datetime.now().strftime("%Y%m%d_%H%M%S")}.docx"'
            response.set_cookie('brief_ready', 'true', max_age=60, path='/')
            return response

    else:
        form = FileFieldForm()

    context= {'form': form,
              'resultados': resultados,
              'preguntas': preguntas,
        }

    return render(request, 'dashboard/dashboard.html',context)

def analizar_archivo(tipo, archivo):
    with tempfile.NamedTemporaryFile(delete=False) as tmp:
        for chunk in archivo.chunks():
            tmp.write(chunk)
        ruta = tmp.name

    texto = extraer_texto(tipo, ruta)
    result_str = extract_insights_with_rag(texto)
    return result_str

def extraer_texto(archivo):
    nombre = archivo.name.lower()
    
    with tempfile.NamedTemporaryFile(delete=False) as tmp:
        for chunk in archivo.chunks():
            tmp.write(chunk)
        ruta = tmp.name

    if nombre.endswith('.pdf'):
        import fitz
        doc = fitz.open(ruta)
        return "\n".join([page.get_text() for page in doc])
    
    elif nombre.endswith('.pptx'):
        from pptx import Presentation
        prs = Presentation(ruta)
        texto = ''
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    texto += shape.text + '\n'
        return texto

    return "No se pudo leer el archivo."

# Define the RAG function
def extract_insights_with_rag(text):
    # 1. Chunk the document
    splitter = RecursiveCharacterTextSplitter(chunk_size=300, chunk_overlap=50)
    docs = [Document(page_content=chunk) for chunk in splitter.split_text(text)]

    # 2. Embeddings
    embeddings = HuggingFaceEmbeddings(model_name="sentence-transformers/all-MiniLM-L6-v2")

    # 3. Vector store
    db = FAISS.from_documents(docs, embeddings)

    # 4. Define retrievera
    retriever = db.as_retriever(search_type="similarity", search_kwargs={"k":3})

    # 5. Use Ollama LLM
    llm = OllamaLLM(model="mistral", format="json")  # o "llama3", si tienes otro modelo cargado

    # 6. Define the prompt/question
    query = """
        You are a strategy analyst assistant. Given the text content of a marketing or business document, extract and return a JSON object containing the following 10 fields.
        Respond ONLY with the JSON object. Do NOT include any additional text, explanations, or markdown fences (json).
        Your entire response must be a single, valid JSON object.
        Fields to extract:

        1. Campaign objectives: What are the business or marketing goals mentioned or implied?
        2. Product specifications: What is being offered? Include relevant product/service details.
        3. Audience segments: Who is the target audience (age, behavior, interests, demographics)?
        4. Content type: What kind of content is being discussed or implied (e.g. blog posts, video, social)?
        5. Tone and style: What tone or voice is preferred (e.g. casual, professional)?
        6. Historical context: Any past efforts, campaigns, or relevant historical background.
        7. Business insights: Any data, findings, or strategic insight present in the document.
        8. Stakeholder_comments: Any input or quotes from internal team members, leadership, etc.
        9. Content tags: Keywords, taxonomy, or themes related to the content.
        10. Delivery platforms: Which channels or platforms are mentioned (e.g. social media, email, website)?

        Respond with the answer in the following format (remember JSON):
    {
        "campaign_objectives": "",
        "product_specifications": "",
        "audience_segments": "",
        "content_type": "",
        "tone_and_style": "",
        "historical_context": "",
        "business_insights": "",
        "stakeholder_comments": "",
        "content_tags": "",
        "delivery_platforms": ""
    }
    """

    # 7. Retrieval-based QA chain
    rag_chain = RetrievalQA.from_chain_type(llm=llm, retriever=retriever, return_source_documents=False)
    result = rag_chain.invoke(query)


    return result

def agentic_objetive(objective: str, model_name: str = "mistral") -> str: #1
    llm = OllamaLLM(model=model_name)

    template = """
    You are an AI agent that receives a business campaign objective and converts it into structured key variables.

    Answer the following in a structured way (DO NOT INFERENCE ANYTHING):
    - Business Objective
    - What are we trying to achieve?
    - Why is this campaign important to the business?
    - Key business KPIs

    Objective:
    {objective}
    """

    prompt = PromptTemplate.from_template(template)
    chain = prompt | llm
    result = chain.invoke({"objective": objective})

    return result

def agentic_product_analysis(product_specifications: str,audience_data: str,content_type: str,tone_style: str, model_name: str = "mistral") -> str: #2
    llm = OllamaLLM(model=model_name)

    template = """
    You are a marketing expert. Analyze the product below to help answer the following (DO NOT INFERENCE ANYTHING):

    1. Marketing Objective
    2. The Problem we are trying to solve
    3. What are the challenges?
    4. Product/Service Description
    5. Solutions/Offering
    6. Why this platform or solution (XYZ)?
    7. Why does the enterprise need it?
    8. 3–5 key features
    9. Main user benefit
    10. What makes it different from competitors?

    ---

    Product description:
    {description}
    Audience:
    {audience_data}
    Type of content:
    {content_type}
    Tone style:
    {tone_style}
    """

    prompt = PromptTemplate.from_template(template)
    chain = prompt | llm
    result = chain.invoke({"description": product_specifications,"audience_data": audience_data,"content_type": content_type,"tone_style": tone_style,})

    return result

def agentic_background(historical_context: str, model_name: str = "mistral") -> str: #3
    llm = OllamaLLM(model=model_name)

    template = """
    You are a meticulous marketing analyst specializing in campaign retrospectives.
    Your task is to carefully review the provided historical campaign data.
    Based on this data, extract and summarize the most critical takeaways (DO NOT INFERENCE ANYTHING).

    Focus on identifying:
    1.  *Key Successes:* What specific aspects or elements of past campaigns worked exceptionally well?
    2.  *Key Failures or Challenges:* What didn't work as expected, what were the main obstacles, or what valuable lessons were learned from setbacks?
    3.  *Notable Creative Themes/Approaches:* Were there any distinct creative styles, messaging tones, or visual approaches that were particularly effective or ineffective?
    4.  *Overall Strategic Implications:* How should these past learnings inform future marketing decisions, strategy adjustments, or avoid repeating mistakes?

    ---

    Past Campaign Data:
    \"\"\"{historical_context}\"\"\"
    ---
    """

    prompt = PromptTemplate.from_template(template)
    chain = prompt | llm
    result = chain.invoke({"historical_context": historical_context})

    return result

def agentic_audience_profiling(audience_data: str,agent_2: str, model_name: str = "mistral") -> str: #4
    llm = OllamaLLM(model=model_name)

    template = """
    You are a digital strategist. Use the audience data to answer (DO NOT INFERENCE ANYTHING):

    1. Target Audience
    2. Ideal Digital Mediums/Channels for this audience
    3. Tone and Style of communication
    4. Key behavioral patterns
    5. Platform preferences

    Audience data:
    {audience_data}
    Product Analysis:
    {product_analysis}
    """

    prompt = PromptTemplate.from_template(template)
    chain = prompt | llm
    result = chain.invoke({"audience_data": audience_data,"product_analysis": agent_2})

    return result

def agentic_insights(product_specifications: str,campaign_objectives: str,agent_4: str,agent_2: str,model_name: str = "mistral") -> str: #5
    llm = OllamaLLM(model=model_name)

    template = """
    You are a strategic planner.

    Analyze the data and provide (DO NOT INFERENCE ANYTHING):

    1. Strategic Insight (linking product + audience + objective)
    2. Creative idea based on that insight
    3. Present market trend and demand
    4. What makes this moment relevant for this campaign

    Product:
    {product}
    Audience:
    {audience}
    Campaign Objective:
    {campaign_objectives}
    Product Analysis:
    {product_analysis}
    """

    prompt = PromptTemplate.from_template(template)
    chain = prompt | llm

    # Pass all three variables to the invoke method
    result = chain.invoke({
        "product": product_specifications,
        "campaign_objectives": campaign_objectives,
        "audience": agent_4,
        "product_analysis": agent_2,
    })

    return result

def agentic_brief_generator(product_specifications: str,agent_4: str,agent_1: str,agent_5: str,agent_3: str,model_name: str = 'mistral') -> str:
    llm = OllamaLLM(model=model_name)

    template = """
    You are a creative strategist at a marketing agency. Build a structured campaign brief with the following sections. Before adding information to the sections, add a line break:

    1. Brief Title
    2. Business Objective
    3. Marketing Objective
    4. Background
    5. Target Audience
    6. The Problem we are trying to solve
    7. What are the challenges?
    8. Product/Service Description
    9. Strategic Insight
    10. Creative Proposal
    11. Present Market Trend and Demand
    12. Digital Channels
    13. Why XYZ (Platform)?
    14. Why does the enterprise need this solution?
    15. Agency Statement of Work (SOW)

    Product:
    {product}

    Audience:
    {audience}

    Objective:
    {objective}

    Insight:
    {insight}

    Background:
    {background}
    """

    prompt = PromptTemplate.from_template(template)
    chain = prompt | llm

    # Invoke the chain with all necessary variables
    result = chain.invoke({
        "product": product_specifications,
        "audience": agent_4,
        "objective": agent_1,
        "insight": agent_5,
        "background": agent_3,
    })

    return result

def generate_brief_docx_in_memory(brief_content: str) -> BytesIO:
    document = DocxDocument()

    # Define the exact list of expected English subtitles for formatting
    EXPECTED_SUBTITLES = [
        "1. Brief Title",
        "2. Business Objective",
        "3. Marketing Objective",
        "4. Background",
        "5. Target Audience",
        "6. The Problem we are trying to solve",
        "7. What are the challenges?",
        "8. Product/Service Description",
        "9. Strategic Insight",
        "10. Creative Proposal",
        "11. Present Market Trend and Demand",
        "12. Digital Channels",
        "13. Why XYZ (Platform)?",
        "14. Why does the enterprise need this solution?",
        "15. Agency Statement of Work (SOW)",
    ]

    # --- Main Title ---
    # Attempt to extract title from the first line if it matches "1. Brief Title"
    main_title_text = "AI-Generated Marketing Campaign Brief" # Default title
    first_line_of_brief = brief_content.split('\n')[0].strip()

    if first_line_of_brief.startswith('1. Brief Title'):
        # Extract title text, removing "1. Brief Title" and any colon
        potential_title = first_line_of_brief.replace('1. Brief Title', '').strip()
        if potential_title.startswith(':'):
            potential_title = potential_title[1:].strip() # Remove leading colon if present
        if potential_title:
            main_title_text = potential_title
    # If the first line doesn't match, the default title "AI-Generated Marketing Campaign Brief" will be used.

    document.add_heading(main_title_text, level=0) # Main title (H1 equivalent)
    document.paragraphs[-1].alignment = WD_ALIGN_PARAGRAPH.CENTER # Center align
    document.add_paragraph() # Add a blank line for spacing

    # --- Process Brief Content Section by Section ---
    lines = brief_content.split('\n')
    
    # Track if we are inside a section to add content as paragraphs
    current_section_found = False

    for line in lines:
        stripped_line = line.strip()
        if not stripped_line: # Ignore empty lines
            continue

        # Check if the current line is one of the expected section titles
        is_section_title = False
        for title_prefix in EXPECTED_SUBTITLES:
            if stripped_line.startswith(title_prefix):
                # Add as a Heading 2
                document.add_heading(stripped_line, level=2)
                is_section_title = True
                current_section_found = True # Mark that we've found the start of a new section
                break # Move to the next line of content

        if not is_section_title:
            # If it's not a new section title, and we've already found at least one section,
            # add the line as a normal paragraph.
            # This handles the actual content belonging to the section.
            if current_section_found:
                document.add_paragraph(stripped_line)
            else:
                # If no section title has been found yet (e.g., if there's leading text before "1. Brief Title"),
                # treat it as a normal paragraph.
                document.add_paragraph(stripped_line)

    # --- Footer ---
    document.add_page_break() # Add a page break at the end
    footer_paragraph = document.add_paragraph("Generated by Briefly - " +
                                              datetime.now().strftime("%Y-%m-%d %H:%M:%S"))
    footer_paragraph.alignment = WD_ALIGN_PARAGRAPH.RIGHT # Right align the footer

    # Save the document to a memory buffer
    buffer = BytesIO()
    document.save(buffer)
    buffer.seek(0) # Rewind the buffer to the beginning
    return buffer

def extract_youtube_text(url):
    try:
        video_id = None
        parsed_url = urlparse(url)
        if parsed_url.hostname == 'youtu.be':
            video_id = parsed_url.path[1:]
        elif parsed_url.hostname in ['www.youtube.com', 'youtube.com']:
            query = parse_qs(parsed_url.query)
            video_id = query.get('v', [None])[0]

        if not video_id:
            return ""

        transcript = YouTubeTranscriptApi.get_transcript(video_id)
        text = " ".join([entry['text'] for entry in transcript])
        return text
    except Exception as e:
        print(f"Error extracting YouTube transcript: {e}")
        return ""
    